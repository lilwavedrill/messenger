# -*- coding: utf-8 -*-
"""Сборка презентации к защите практики.
Стиль: минимализм в светлой теме, единственный акцент — электрик-синий.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

# --- палитра и типографика ---
BG      = RGBColor(0xFF, 0xFF, 0xFF)
INK     = RGBColor(0x0A, 0x0A, 0x0A)  # почти чёрный
MUTED   = RGBColor(0x52, 0x52, 0x5B)
HAIR    = RGBColor(0xE5, 0xE5, 0xE5)
FAINT   = RGBColor(0xF5, 0xF5, 0xF5)
BLUE    = RGBColor(0x00, 0x66, 0xFF)
BLUE_BG = RGBColor(0xF5, 0xF8, 0xFF)  # почти-белый с синевой
GREEN   = RGBColor(0x22, 0xC5, 0x5E)

# Шрифты — стандартные Windows/Office, без зависимости от Google Fonts
F_UI    = "Segoe UI"
F_MONO  = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

DOCS = Path(__file__).parent


# ---------- utilities ----------

def add_blank_slide(prs):
    blank = prs.slide_layouts[6]  # blank
    s = prs.slides.add_slide(blank)
    # белый фон
    bg = s.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    return s


def add_text(slide, left, top, width, height, text, *,
             font=F_UI, size=18, bold=False, italic=False,
             color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             letter_spacing=0):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def add_multiline(slide, left, top, width, height, lines, *,
                  font=F_UI, size=16, color=INK, bold=False,
                  line_spacing=1.35, bullet_dash=True):
    """Многострочный текст с одинаковым стилем; bullets через тире."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        r = p.add_run()
        prefix = "— " if bullet_dash else ""
        r.text = prefix + line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def add_rect(slide, left, top, width, height, fill, line=None, line_w=0):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def hairline(slide, x, y, w):
    add_rect(slide, x, y, w, Emu(9525), HAIR)  # ~0.75pt высоты


def add_dot_grid(slide):
    """Маленькая 2×2 сетка точек — фирменный знак."""
    x, y, sz, gap = Inches(0.55), Inches(0.5), Emu(160000), Emu(70000)
    for r in range(2):
        for c in range(2):
            add_rect(slide,
                     x + c * (sz + gap),
                     y + r * (sz + gap),
                     sz, sz, BLUE)


def add_slide_header(slide, kicker, title):
    """Заголовок слайда: моно-килька сверху, крупный заголовок под ним, hairline."""
    add_dot_grid(slide)
    add_text(slide, Inches(1.1), Inches(0.42), Inches(6), Inches(0.35),
             kicker, font=F_MONO, size=11, color=MUTED)
    add_text(slide, Inches(0.55), Inches(0.95), Inches(12), Inches(0.85),
             title, font=F_UI, size=32, bold=True, color=INK)
    hairline(slide, Inches(0.55), Inches(1.85), Inches(12.2))


def add_footer(slide, num, total):
    """Тонкий футер: слева бренд, справа номер слайда."""
    add_text(slide, Inches(0.55), Inches(7.05), Inches(6), Inches(0.3),
             "MESSENGER  ·  VK Education / Case #1",
             font=F_MONO, size=9, color=MUTED)
    add_text(slide, Inches(11.55), Inches(7.05), Inches(1.7), Inches(0.3),
             f"{num:02d} / {total:02d}",
             font=F_MONO, size=9, color=MUTED, align=PP_ALIGN.RIGHT)


# ---------- слайды ----------

def slide_title(prs):
    s = add_blank_slide(prs)
    # синяя вертикальная полоса слева
    add_rect(s, 0, 0, Inches(0.35), SLIDE_H, BLUE)
    # dot-grid как визуальный акцент
    x, y, sz, gap = Inches(1.2), Inches(1.4), Emu(360000), Emu(140000)
    for r in range(2):
        for c in range(2):
            add_rect(s, x + c * (sz + gap), y + r * (sz + gap), sz, sz, BLUE)
    # kicker
    add_text(s, Inches(1.2), Inches(3.1), Inches(11), Inches(0.4),
             "ОТЧЁТ ПО ПРАКТИКЕ  ·  КЕЙС №1  ·  VK EDUCATION",
             font=F_MONO, size=12, color=MUTED)
    # заголовок
    add_text(s, Inches(1.2), Inches(3.55), Inches(11), Inches(1.2),
             "Мессенджер",
             font=F_UI, size=72, bold=True, color=INK)
    # описание
    add_text(s, Inches(1.2), Inches(4.9), Inches(11), Inches(0.5),
             "Реалтайм-обмен сообщениями, вложения, полнотекстовый поиск.",
             font=F_UI, size=18, color=MUTED)
    # авторы + вуз (правый нижний угол)
    add_text(s, Inches(1.2), Inches(6.05), Inches(11), Inches(0.35),
             "НИТУ МИСИС  ·  группа БИВТ-24-1",
             font=F_MONO, size=11, color=MUTED)
    add_multiline(s, Inches(1.2), Inches(6.35), Inches(11), Inches(0.8), [
        "Мамонтов Иван Андреевич",
        "Кривонос Никита",
        "Руководитель практики: Чумакова М. Ю.",
    ], size=13, color=INK, line_spacing=1.15, bullet_dash=False)


def slide_agenda(prs, total):
    s = add_blank_slide(prs)
    add_slide_header(s, "ПЛАН", "О чём расскажем")
    items = [
        "Актуальность и задачи",
        "Требования и обзор аналогов",
        "Стек технологий",
        "Архитектура системы",
        "База данных и полнотекстовый поиск",
        "Модель авторизации",
        "Интерфейс приложения",
        "Тестирование и результаты",
    ]
    # две колонки
    col_w = Inches(5.6)
    left1 = Inches(1.0)
    left2 = Inches(7.2)
    mid = len(items) // 2 + len(items) % 2
    add_multiline(s, left1, Inches(2.4), col_w, Inches(4.5), items[:mid],
                  size=20, color=INK, line_spacing=1.6)
    add_multiline(s, left2, Inches(2.4), col_w, Inches(4.5), items[mid:],
                  size=20, color=INK, line_spacing=1.6)
    add_footer(s, 2, total)


def slide_relevance(prs, total):
    s = add_blank_slide(prs)
    add_slide_header(s, "01  ·  АКТУАЛЬНОСТЬ", "Почему это интересно")

    # три цифры-хайлайта
    stats = [
        ("100+ млрд", "сообщений в день через мессенджеры"),
        ("< 200 мс", "цель по задержке доставки"),
        ("3 транспорта", "WebSocket, gRPC, long polling"),
    ]
    x = Inches(0.9)
    y = Inches(2.4)
    card_w = Inches(3.9)
    card_h = Inches(1.8)
    gap = Inches(0.25)
    for i, (num, label) in enumerate(stats):
        left = x + i * (card_w + gap)
        add_rect(s, left, y, card_w, card_h, FAINT)
        add_rect(s, left, y, Emu(35000), card_h, BLUE)  # тонкая синяя полоса слева
        add_text(s, left + Inches(0.35), y + Inches(0.28),
                 card_w - Inches(0.4), Inches(0.6),
                 num, font=F_UI, size=28, bold=True, color=INK)
        add_text(s, left + Inches(0.35), y + Inches(0.95),
                 card_w - Inches(0.4), Inches(0.7),
                 label, font=F_UI, size=13, color=MUTED)

    add_multiline(s, Inches(0.9), Inches(4.7), Inches(11.5), Inches(2), [
        "За внешней простотой чата стоят нетривиальные задачи: минимальная задержка,",
        "устойчивость к обрывам сети, защита от дубликатов, быстрый поиск по истории.",
        "Разработка учебного мессенджера позволяет разобрать эти задачи на практике.",
    ], size=15, color=INK, line_spacing=1.5, bullet_dash=False)
    add_footer(s, 3, total)


def slide_goals(prs, total):
    s = add_blank_slide(prs)
    add_slide_header(s, "02  ·  ЦЕЛЬ И ЗАДАЧИ", "Что должно получиться")

    # цель
    add_text(s, Inches(0.9), Inches(2.35), Inches(11.5), Inches(0.35),
             "ЦЕЛЬ", font=F_MONO, size=11, color=MUTED)
    add_multiline(s, Inches(0.9), Inches(2.7), Inches(11.5), Inches(0.9), [
        "Спроектировать и реализовать рабочий прототип мессенджера, отвечающий",
        "функциональным и нефункциональным требованиям технического задания.",
    ], size=17, color=INK, line_spacing=1.35, bullet_dash=False)

    add_text(s, Inches(0.9), Inches(4.05), Inches(11.5), Inches(0.35),
             "ЗАДАЧИ", font=F_MONO, size=11, color=MUTED)
    tasks = [
        "проанализировать предметную область и требования ТЗ",
        "спроектировать схему БД и API",
        "реализовать серверную часть на Python + FastAPI",
        "организовать реалтайм-обмен через WebSocket",
        "построить авторизацию на JWT + refresh с ротацией",
        "разработать веб-клиент",
        "покрыть основные сценарии автотестами",
    ]
    add_multiline(s, Inches(0.9), Inches(4.45), Inches(11.5), Inches(2.5),
                  tasks, size=15, color=INK, line_spacing=1.35)
    add_footer(s, 4, total)


def slide_requirements(prs, total):
    s = add_blank_slide(prs)
    add_slide_header(s, "03  ·  ТРЕБОВАНИЯ", "11 функциональных + 6 нефункциональных")

    # две колонки: обязательное / дополнительное
    y0 = Inches(2.4)
    add_text(s, Inches(0.9), y0, Inches(6), Inches(0.35),
             "ОБЯЗАТЕЛЬНЫЕ", font=F_MONO, size=11, color=BLUE)
    add_multiline(s, Inches(0.9), y0 + Inches(0.35), Inches(6), Inches(3.5), [
        "регистрация и вход",
        "личные и групповые чаты",
        "реалтайм-обмен сообщениями",
        "история с пагинацией",
        "управление участниками (owner)",
        "поиск по сообщениям чата",
        "разграничение ролей внутри чата",
    ], size=14, color=INK, line_spacing=1.5)

    add_text(s, Inches(7.2), y0, Inches(6), Inches(0.35),
             "ДОПОЛНИТЕЛЬНО", font=F_MONO, size=11, color=BLUE)
    add_multiline(s, Inches(7.2), y0 + Inches(0.35), Inches(6), Inches(3.5), [
        "статусы доставки (sent / delivered / read)",
        "идемпотентность через client_msg_id",
        "вложения через S3-совместимое хранилище",
        "refresh-токены JWT с ротацией",
        "аудит действий пользователей",
    ], size=14, color=INK, line_spacing=1.5)

    # НФТ снизу — плашкой
    add_rect(s, Inches(0.9), Inches(6.15), Inches(12.3), Inches(0.7), FAINT)
    add_text(s, Inches(1.1), Inches(6.22), Inches(12), Inches(0.3),
             "НЕФУНКЦИОНАЛЬНЫЕ", font=F_MONO, size=10, color=MUTED)
    add_text(s, Inches(1.1), Inches(6.5), Inches(12), Inches(0.3),
             "отзывчивость < 200 мс  ·  надёжность  ·  безопасность  ·  порядок  ·  наблюдаемость  ·  совместимость",
             font=F_UI, size=12, color=INK)
    add_footer(s, 5, total)


def slide_analogs(prs, total):
    s = add_blank_slide(prs)
    add_slide_header(s, "04  ·  ОБЗОР АНАЛОГОВ", "Что делают взрослые, а что мы")

    items = [
        ("WhatsApp",  "XMPP-подобный протокол, Signal E2E-шифрование"),
        ("Telegram",  "собственный MTProto, распределённые дата-центры"),
        ("Discord",   "интенсивно использует WebSocket, ориентация на голос"),
        ("Slack",     "гибрид REST + WebSocket — та же схема, что у нас"),
        ("VK",        "long polling в мобильных, WebSocket в web"),
        ("Matrix",    "открытый федеративный протокол"),
    ]
    y = Inches(2.35)
    for i, (name, note) in enumerate(items):
        row_y = y + Inches(0.55) * i
        add_text(s, Inches(0.9), row_y, Inches(2.5), Inches(0.5),
                 name, font=F_UI, size=16, bold=True, color=INK)
        add_text(s, Inches(3.6), row_y, Inches(9.6), Inches(0.5),
                 note, font=F_UI, size=14, color=MUTED)

    # что взяли / чего избежали
    add_rect(s, Inches(0.9), Inches(6.0), Inches(6.0), Inches(0.9), BLUE_BG)
    add_text(s, Inches(1.1), Inches(6.08), Inches(6), Inches(0.3),
             "ВЗЯЛИ ЗА ОСНОВУ", font=F_MONO, size=9, color=BLUE)
    add_text(s, Inches(1.1), Inches(6.36), Inches(5.8), Inches(0.5),
             "гибрид REST + WebSocket, роли в чате, идемпотентность",
             font=F_UI, size=12, color=INK)

    add_rect(s, Inches(7.2), Inches(6.0), Inches(6.0), Inches(0.9), FAINT)
    add_text(s, Inches(7.4), Inches(6.08), Inches(6), Inches(0.3),
             "СОЗНАТЕЛЬНО ОТКАЗАЛИСЬ", font=F_MONO, size=9, color=MUTED)
    add_text(s, Inches(7.4), Inches(6.36), Inches(5.8), Inches(0.5),
             "E2E-шифрование, федерация, брокер сообщений, шардирование",
             font=F_UI, size=12, color=INK)
    add_footer(s, 6, total)


def slide_stack(prs, total):
    s = add_blank_slide(prs)
    add_slide_header(s, "05  ·  СТЕК", "Технологии и почему")

    items = [
        ("Python 3.13",   "серверный язык"),
        ("FastAPI",       "REST + WebSocket в одном процессе, авто-OpenAPI"),
        ("SQLAlchemy 2",  "ORM, синхронный режим"),
        ("PostgreSQL 17", "транзакции, JSONB, GIN-индексы"),
        ("MinIO (S3)",    "объектное хранилище для вложений"),
        ("bcrypt + JWT",  "хэши паролей, access + refresh токены"),
        ("pytest + httpx", "автотесты API"),
        ("HTML + CSS + JS", "клиент без сборщика и без фреймворков"),
    ]
    # сетка 2×4
    x0 = Inches(0.9)
    y0 = Inches(2.35)
    card_w = Inches(6.1)
    card_h = Inches(1.05)
    gap_x = Inches(0.15)
    gap_y = Inches(0.15)
    for i, (name, desc) in enumerate(items):
        col = i % 2
        row = i // 2
        left = x0 + col * (card_w + gap_x)
        top = y0 + row * (card_h + gap_y)
        add_rect(s, left, top, card_w, card_h, FAINT)
        add_rect(s, left, top, Emu(35000), card_h, BLUE)
        add_text(s, left + Inches(0.3), top + Inches(0.15),
                 card_w - Inches(0.4), Inches(0.4),
                 name, font=F_UI, size=14, bold=True, color=INK)
        add_text(s, left + Inches(0.3), top + Inches(0.55),
                 card_w - Inches(0.4), Inches(0.4),
                 desc, font=F_UI, size=11, color=MUTED)
    add_footer(s, 7, total)


def slide_architecture(prs, total):
    s = add_blank_slide(prs)
    add_slide_header(s, "06  ·  АРХИТЕКТУРА", "Трёхзвенная схема + S3")

    # блоки: клиент — сервер — БД, S3 сбоку
    def block(left, top, w, h, title_top, title_bottom, fill=FAINT):
        add_rect(s, left, top, w, h, fill, line=HAIR, line_w=0.75)
        add_text(s, left, top + Inches(0.35), w, Inches(0.4),
                 title_top, font=F_UI, size=16, bold=True, color=INK,
                 align=PP_ALIGN.CENTER)
        add_text(s, left, top + Inches(0.75), w, Inches(0.35),
                 title_bottom, font=F_MONO, size=11, color=MUTED,
                 align=PP_ALIGN.CENTER)

    y = Inches(2.7)
    w = Inches(3.2)
    h = Inches(1.4)
    # клиент
    block(Inches(0.9), y, w, h, "Клиент", "HTML + CSS + JS")
    # сервер
    block(Inches(5.05), y, w, h, "Сервер приложения", "FastAPI + WS Hub", fill=BLUE_BG)
    # БД
    block(Inches(9.2), y, w, h, "PostgreSQL", "8 таблиц + GIN FTS")
    # стрелки текстом
    add_text(s, Inches(4.1), y + Inches(0.5), Inches(0.95), Inches(0.4),
             "REST / WS", font=F_MONO, size=11, color=MUTED,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(8.25), y + Inches(0.5), Inches(0.95), Inches(0.4),
             "SQLAlchemy", font=F_MONO, size=11, color=MUTED,
             align=PP_ALIGN.CENTER)
    # S3 внизу
    block(Inches(5.05), Inches(4.6), w, h, "MinIO (S3)", "вложения + аватары")
    add_text(s, Inches(6.15), Inches(4.25), Inches(1), Inches(0.3),
             "boto3", font=F_MONO, size=11, color=MUTED,
             align=PP_ALIGN.CENTER)

    add_multiline(s, Inches(0.9), Inches(6.35), Inches(12.3), Inches(0.9), [
        "Гибрид: сообщение отправляется REST-запросом (транзакция в БД), затем",
        "рассылается через WebSocket подписчикам чата. Даже при обрыве WS ничего не теряется.",
    ], size=13, color=INK, line_spacing=1.35, bullet_dash=False)
    add_footer(s, 8, total)


def slide_database(prs, total):
    s = add_blank_slide(prs)
    add_slide_header(s, "07  ·  БАЗА ДАННЫХ", "8 таблиц, 3НФ, индексы под задачу")

    # список таблиц слева
    entities = [
        "users",
        "chats",
        "chat_members",
        "messages",
        "message_status",
        "attachments",
        "refresh_tokens",
        "audit_log",
    ]
    add_text(s, Inches(0.9), Inches(2.35), Inches(4), Inches(0.35),
             "ТАБЛИЦЫ", font=F_MONO, size=11, color=MUTED)
    add_multiline(s, Inches(0.9), Inches(2.75), Inches(4), Inches(3.5),
                  entities, size=15, color=INK, line_spacing=1.45,
                  font=F_MONO)

    # ключевые индексы справа
    add_text(s, Inches(5.5), Inches(2.35), Inches(7.8), Inches(0.35),
             "КЛЮЧЕВЫЕ ИНДЕКСЫ", font=F_MONO, size=11, color=MUTED)
    idx = [
        "(chat_id, id DESC) на messages — пагинация «вверх» O(log n)",
        "GIN to_tsvector('russian', text) на messages — FTS с русской морфологией",
        "(user_id) на chat_members — быстрая выборка «мои чаты»",
        "UNIQUE (chat_id, client_msg_id) — механизм идемпотентности",
        "UNIQUE token_hash на refresh_tokens",
    ]
    add_multiline(s, Inches(5.5), Inches(2.75), Inches(7.8), Inches(4),
                  idx, size=13, color=INK, line_spacing=1.45)

    # ссылка на приложение
    add_text(s, Inches(0.9), Inches(6.55), Inches(12.3), Inches(0.35),
             "ER-диаграмма — в приложении А отчёта (docs/er.png).",
             font=F_UI, size=12, italic=True, color=MUTED)
    add_footer(s, 9, total)


def slide_fts(prs, total):
    s = add_blank_slide(prs)
    add_slide_header(s, "08  ·  ПОЛНОТЕКСТОВЫЙ ПОИСК", "От LIKE к GIN + русскому словарю")

    # было
    add_rect(s, Inches(0.9), Inches(2.35), Inches(5.9), Inches(2.1), FAINT)
    add_text(s, Inches(1.1), Inches(2.5), Inches(5.6), Inches(0.35),
             "БЫЛО", font=F_MONO, size=11, color=MUTED)
    add_text(s, Inches(1.1), Inches(2.85), Inches(5.6), Inches(0.5),
             "WHERE LOWER(text) LIKE '%q%'",
             font=F_MONO, size=13, color=INK)
    add_multiline(s, Inches(1.1), Inches(3.4), Inches(5.6), Inches(1.1), [
        "полный скан таблицы",
        "не находит другие словоформы",
    ], size=12, color=MUTED, line_spacing=1.3)

    # стало
    add_rect(s, Inches(7.0), Inches(2.35), Inches(5.9), Inches(2.1), BLUE_BG)
    add_rect(s, Inches(7.0), Inches(2.35), Emu(35000), Inches(2.1), BLUE)
    add_text(s, Inches(7.2), Inches(2.5), Inches(5.6), Inches(0.35),
             "СТАЛО", font=F_MONO, size=11, color=BLUE)
    add_text(s, Inches(7.2), Inches(2.85), Inches(5.6), Inches(0.5),
             "to_tsvector('russian') @@ plainto_tsquery",
             font=F_MONO, size=12, color=INK)
    add_multiline(s, Inches(7.2), Inches(3.4), Inches(5.6), Inches(1.1), [
        "GIN-индекс, никакого скана",
        "стемминг: «обсуждать», «обсуждаю», «обсуждали»",
        "ранжирование через ts_rank_cd",
    ], size=12, color=INK, line_spacing=1.3)

    # пример SQL
    add_text(s, Inches(0.9), Inches(4.85), Inches(12), Inches(0.35),
             "ЗАПРОС", font=F_MONO, size=11, color=MUTED)
    add_rect(s, Inches(0.9), Inches(5.15), Inches(12.3), Inches(1.5), FAINT)
    add_multiline(s, Inches(1.05), Inches(5.25), Inches(12), Inches(1.4), [
        "SELECT *,",
        "       ts_rank_cd(to_tsvector('russian', text), plainto_tsquery('russian', :q)) AS rank",
        "  FROM messages",
        " WHERE chat_id = :chat_id",
        "   AND to_tsvector('russian', text) @@ plainto_tsquery('russian', :q)",
        " ORDER BY rank DESC, id DESC;",
    ], size=11, color=INK, line_spacing=1.15, bullet_dash=False, font=F_MONO)
    add_footer(s, 10, total)


def slide_auth(prs, total):
    s = add_blank_slide(prs)
    add_slide_header(s, "09  ·  АВТОРИЗАЦИЯ", "JWT + refresh + ротация")

    # схема
    stages = [
        ("POST /auth/login",   "пара access + refresh"),
        ("Bearer <access>",    "все защищённые запросы, TTL 30 мин"),
        ("POST /auth/refresh", "старый refresh отзывается, выдаётся новая пара"),
        ("POST /auth/logout",  "refresh помечается revoked"),
    ]
    y = Inches(2.35)
    for i, (endpoint, note) in enumerate(stages):
        row_y = y + Inches(0.85) * i
        add_rect(s, Inches(0.9), row_y, Inches(0.35), Inches(0.6), BLUE)
        add_text(s, Inches(1.4), row_y + Inches(0.05), Inches(5),
                 Inches(0.4), endpoint,
                 font=F_MONO, size=13, bold=True, color=INK)
        add_text(s, Inches(1.4), row_y + Inches(0.35), Inches(11), Inches(0.4),
                 note, font=F_UI, size=13, color=MUTED)

    # блок безопасности справа
    add_rect(s, Inches(7.6), Inches(2.35), Inches(5.6), Inches(3.4), FAINT)
    add_text(s, Inches(7.8), Inches(2.5), Inches(5.2), Inches(0.35),
             "БЕЗОПАСНОСТЬ", font=F_MONO, size=11, color=MUTED)
    add_multiline(s, Inches(7.8), Inches(2.85), Inches(5.2), Inches(2.9), [
        "пароли — только bcrypt-хэш",
        "refresh хранится SHA-256 хэшем",
        "ротация: одноразовый refresh",
        "row-lock (SELECT FOR UPDATE) от гонки",
        "revoked_at помечает отозванные",
        "TTL access = 30 мин · refresh = 30 дней",
    ], size=13, color=INK, line_spacing=1.5)
    add_footer(s, 11, total)


def slide_screenshots(prs, total):
    s = add_blank_slide(prs)
    add_slide_header(s, "10  ·  ИНТЕРФЕЙС", "Скриншоты работающего клиента")

    shots = [
        ("screenshots/01-login.png",    "Форма входа"),
        ("screenshots/02-main.png",     "Основной экран"),
        ("screenshots/03-profile.png",  "Профиль"),
        ("screenshots/04-new-chat.png", "Создание группы"),
    ]
    # 2×2 сетка изображений
    x0 = Inches(0.9)
    y0 = Inches(2.3)
    cell_w = Inches(6.1)
    cell_h = Inches(2.1)
    gap_x = Inches(0.15)
    gap_y = Inches(0.35)
    for i, (rel, caption) in enumerate(shots):
        col = i % 2
        row = i // 2
        left = x0 + col * (cell_w + gap_x)
        top = y0 + row * (cell_h + gap_y)
        # фоновая плашка + рамка hairline
        add_rect(s, left, top, cell_w, cell_h, FAINT, line=HAIR, line_w=0.75)
        # картинка по центру ячейки — python-pptx масштабирует по width, aspect сохраняется
        path = DOCS / rel
        if path.exists():
            # оценим ширину, пусть будет чуть меньше ячейки, чтобы влезло
            pic = s.shapes.add_picture(str(path),
                                       left + Inches(0.25),
                                       top + Inches(0.15),
                                       width=cell_w - Inches(0.5))
            # если высота вышла за пределы — пересчитаем через height
            if pic.height > cell_h - Inches(0.55):
                sp = pic._element
                sp.getparent().remove(sp)
                pic = s.shapes.add_picture(str(path),
                                           left + Inches(0.25),
                                           top + Inches(0.15),
                                           height=cell_h - Inches(0.55))
                # центрирование по горизонтали
                pic.left = left + (cell_w - pic.width) // 2
        # подпись под ячейкой
        add_text(s, left, top + cell_h - Inches(0.35), cell_w, Inches(0.3),
                 caption, font=F_MONO, size=10, color=MUTED,
                 align=PP_ALIGN.CENTER)
    add_footer(s, 12, total)


def slide_testing(prs, total):
    s = add_blank_slide(prs)
    add_slide_header(s, "11  ·  ТЕСТИРОВАНИЕ", "Автотесты + ручные сценарии")

    # цифры-хайлайты
    stats = [
        ("7 / 7",  "pytest-тестов проходят"),
        ("4.5 сек", "полный прогон suite"),
        ("6 сцен.", "ручных end-to-end проверок"),
    ]
    x = Inches(0.9)
    y = Inches(2.35)
    card_w = Inches(3.9)
    card_h = Inches(1.7)
    gap = Inches(0.25)
    for i, (num, label) in enumerate(stats):
        left = x + i * (card_w + gap)
        add_rect(s, left, y, card_w, card_h, BLUE_BG)
        add_rect(s, left, y, Emu(35000), card_h, BLUE)
        add_text(s, left + Inches(0.35), y + Inches(0.25),
                 card_w - Inches(0.4), Inches(0.6),
                 num, font=F_UI, size=26, bold=True, color=INK)
        add_text(s, left + Inches(0.35), y + Inches(0.9),
                 card_w - Inches(0.4), Inches(0.7),
                 label, font=F_UI, size=13, color=MUTED)

    # что покрыто
    add_text(s, Inches(0.9), Inches(4.45), Inches(12), Inches(0.35),
             "ЧТО ПОКРЫТО", font=F_MONO, size=11, color=MUTED)
    covered = [
        "регистрация и вход по логину/паролю",
        "неверный пароль → 401",
        "создание чата, отправка, история, поиск (со стеммингом)",
        "ограничение прав: reader не может отправлять — 403",
        "ротация refresh: старый refresh становится недействительным",
        "идемпотентность: повтор с тем же client_msg_id не создаёт дубликат",
    ]
    add_multiline(s, Inches(0.9), Inches(4.85), Inches(12), Inches(2.1),
                  covered, size=13, color=INK, line_spacing=1.35)
    add_footer(s, 13, total)


def slide_results(prs, total):
    s = add_blank_slide(prs)
    add_slide_header(s, "12  ·  РЕЗУЛЬТАТЫ", "Что готово и что дальше")

    add_text(s, Inches(0.9), Inches(2.35), Inches(12), Inches(0.35),
             "ГОТОВО", font=F_MONO, size=11, color=BLUE)
    done = [
        "все 7 обязательных функциональных требований",
        "4 дополнительных (статусы, вложения, refresh, идемпотентность)",
        "полнотекстовый поиск на GIN + русской морфологии",
        "автоматизированные тесты (7 / 7 проходят)",
        "отчёт по ГОСТ 7.32-2017, ER-диаграмма, видеодемонстрация",
    ]
    add_multiline(s, Inches(0.9), Inches(2.75), Inches(12), Inches(2.5),
                  done, size=14, color=INK, line_spacing=1.45)

    add_text(s, Inches(0.9), Inches(5.05), Inches(12), Inches(0.35),
             "НАПРАВЛЕНИЯ РАЗВИТИЯ", font=F_MONO, size=11, color=MUTED)
    todo = [
        "внешний брокер (Redis Streams / NATS) для горизонтального масштабирования",
        "оффлайн-уведомления, редактирование и удаление сообщений",
        "E2E-шифрование, голосовые сообщения, WebRTC-видеозвонки",
    ]
    add_multiline(s, Inches(0.9), Inches(5.45), Inches(12), Inches(1.6),
                  todo, size=14, color=INK, line_spacing=1.45)

    # ссылка на репозиторий
    add_rect(s, Inches(0.9), Inches(6.9), Inches(12.3), Emu(9525), HAIR)
    add_text(s, Inches(0.9), Inches(7.0), Inches(12.3), Inches(0.3),
             "github.com/lilwavedrill/messenger",
             font=F_MONO, size=11, color=BLUE)


# ---------- main ----------

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    total = 13  # 1 титул + 12 контентных

    slide_title(prs)
    slide_agenda(prs, total)
    slide_relevance(prs, total)
    slide_goals(prs, total)
    slide_requirements(prs, total)
    slide_analogs(prs, total)
    slide_stack(prs, total)
    slide_architecture(prs, total)
    slide_database(prs, total)
    slide_fts(prs, total)
    slide_auth(prs, total)
    slide_screenshots(prs, total)
    slide_testing(prs, total)
    slide_results(prs, total)

    out = DOCS / "presentation.pptx"
    prs.save(str(out))
    print(f"OK, wrote {out} ({out.stat().st_size // 1024} KB)")


if __name__ == "__main__":
    main()
